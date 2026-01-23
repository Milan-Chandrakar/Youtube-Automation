"""
Report generation module for creating professional presentations.
Uses Plotly for visualizations and python-pptx for slide deck creation.
"""

import logging
from datetime import datetime
from pathlib import Path
from typing import Dict, List
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from config import REPORT_OUTPUT_PATH, REPORT_TITLE

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ReportGenerator:
    """Generates professional presentation reports with visualizations."""

    def __init__(self, insights: Dict, df: pd.DataFrame):
        self.insights = insights
        self.df = df
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.charts_dir = REPORT_OUTPUT_PATH / "charts"
        self.charts_dir.mkdir(parents=True, exist_ok=True)

    def create_title_slide(self):
        """Create title slide."""
        logger.info("Creating title slide")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(44, 62, 80)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = REPORT_TITLE
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(189, 195, 199)
        p.alignment = PP_ALIGN.CENTER

    def create_executive_summary_slide(self):
        """Create executive summary slide."""
        logger.info("Creating executive summary slide")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank layout
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)

        # Summary statistics
        stats = self.insights.get("engagement_stats", {})
        sentiment = self.insights.get("sentiment_distribution", {})

        summary_text = f"""
Videos Analyzed: {len(self.df)}

Engagement Metrics:
• Average Views per Video: {stats.get('avg_views', 0):,.0f}
• Average Engagement Rate: {stats.get('engagement_rate', 0):.2%}
• Total Comments: {self.df['comment_count'].sum():,}

Sentiment Analysis:
• Positive Sentiment: {sentiment.get('positive', 0)} videos
• Negative Sentiment: {sentiment.get('negative', 0)} videos
• Neutral Sentiment: {sentiment.get('neutral', 0)} videos
• Average Sentiment Score: {sentiment.get('avg_sentiment', 0):.2f}
        """

        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = summary_text.strip()
        p.font.size = Pt(16)
        p.line_spacing = 1.5

    def create_engagement_chart(self):
        """Create engagement metrics chart."""
        logger.info("Creating engagement chart")

        try:
            # Top videos by engagement
            top_videos = self.df.nlargest(10, "engagement_score")[
                ["title", "engagement_score", "view_count"]
            ].copy()

            fig = px.bar(
                top_videos,
                x="engagement_score",
                y="title",
                orientation="h",
                title="Top 10 Videos by Engagement Score",
                labels={"engagement_score": "Engagement Score", "title": "Video Title"},
            )
            fig.update_layout(height=600, showlegend=False, yaxis_tickfont=dict(size=10))

            chart_path = self.charts_dir / "engagement_chart.png"
            try:
                fig.write_image(chart_path, width=1200, height=600)
                logger.info(f"Saved chart: {chart_path}")
            except Exception as e:
                logger.warning(f"Could not save chart image: {e}. Continuing without image.")
                return None

            return chart_path
        except Exception as e:
            logger.warning(f"Error creating engagement chart: {e}")
            return None

    def create_sentiment_chart(self):
        """Create sentiment distribution chart."""
        logger.info("Creating sentiment chart")

        try:
            sentiment = self.insights.get("sentiment_distribution", {})
            data = {
                "Sentiment": ["Positive", "Negative", "Neutral"],
                "Count": [
                    sentiment.get("positive", 0),
                    sentiment.get("negative", 0),
                    sentiment.get("neutral", 0),
                ],
            }

            fig = go.Figure(
                data=[
                    go.Pie(
                        labels=data["Sentiment"],
                        values=data["Count"],
                        marker=dict(colors=["#27ae60", "#e74c3c", "#95a5a6"]),
                    )
                ]
            )
            fig.update_layout(title="Sentiment Distribution", height=500)

            chart_path = self.charts_dir / "sentiment_chart.png"
            try:
                fig.write_image(chart_path, width=800, height=600)
                logger.info(f"Saved chart: {chart_path}")
            except Exception as e:
                logger.warning(f"Could not save sentiment chart image: {e}. Continuing without image.")
                return None

            return chart_path
        except Exception as e:
            logger.warning(f"Error creating sentiment chart: {e}")
            return None

    def create_themes_slide(self):
        """Create industry themes slide."""
        logger.info("Creating themes slide")

        themes = self.insights.get("industry_themes", {}).get("top_themes", [])

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Key Industry Themes"
        p.font.size = Pt(44)
        p.font.bold = True

        themes_text = "\n".join([f"• {theme}" for theme in themes[:12]])

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = themes_text
        p.font.size = Pt(18)
        p.line_spacing = 1.8

    def add_chart_slide(self, chart_title: str, chart_path: Path):
        """Add a slide with a chart image."""
        if chart_path is None:
            logger.warning(f"Skipping chart slide: {chart_title} (no image available)")
            return
            
        logger.info(f"Adding chart slide: {chart_title}")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = chart_title
        p.font.size = Pt(40)
        p.font.bold = True

        # Add chart image
        if chart_path.exists():
            slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.2), width=Inches(9))

    def create_top_channels_slide(self):
        """Create top channels slide."""
        logger.info("Creating top channels slide")

        channels = self.insights.get("top_channels", {})

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Top Performing Channels"
        p.font.size = Pt(40)
        p.font.bold = True

        channels_text = ""
        for i, (channel, stats) in enumerate(list(channels.items())[:8], 1):
            channels_text += f"{i}. {channel}\n"
            channels_text += f"   Videos: {stats.get('video_count', 0)}, Total Views: {stats.get('total_views', 0):,}\n\n"

        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = channels_text.strip()
        p.font.size = Pt(14)

    def create_conclusion_slide(self):
        """Create conclusion slide."""
        logger.info("Creating conclusion slide")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Key Takeaways"
        p.font.size = Pt(40)
        p.font.bold = True

        conclusion = """
• The AI and automation niche shows strong engagement and growth
• Sentiment across the industry is predominantly positive
• Top creators focus on practical applications and tutorials
• Video length and production quality correlate with engagement
• Trending topics revolve around LLMs, automation workflows, and implementation guides

Recommendations:
✓ Focus on in-depth, practical content
✓ Engage with trending tools and frameworks
✓ Build community through comments and interactions
✓ Maintain consistent upload schedule
✓ Optimize titles and descriptions for discoverability
        """

        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = conclusion.strip()
        p.font.size = Pt(14)
        p.line_spacing = 1.6

    def generate_presentation(self, output_filename: str = "YouTube_Trends_Report.pptx") -> Path:
        """
        Generate complete presentation with all slides and charts.

        Args:
            output_filename: Name of output PowerPoint file

        Returns:
            Path to generated presentation
        """
        logger.info("Generating presentation")

        self.create_title_slide()
        self.create_executive_summary_slide()

        engagement_chart = self.create_engagement_chart()
        self.add_chart_slide("Engagement Analysis", engagement_chart)

        sentiment_chart = self.create_sentiment_chart()
        self.add_chart_slide("Sentiment Distribution", sentiment_chart)

        self.create_themes_slide()
        self.create_top_channels_slide()
        self.create_conclusion_slide()

        output_path = REPORT_OUTPUT_PATH / output_filename
        self.prs.save(str(output_path))
        logger.info(f"Presentation saved to {output_path}")

        return output_path


def main():
    """Example usage."""
    from analyzer import YouTubeAnalyzer

    df = pd.read_csv("data/videos.csv")
    analyzer = YouTubeAnalyzer(df)
    insights = analyzer.run_full_analysis()

    generator = ReportGenerator(insights, df)
    report_path = generator.generate_presentation()
    print(f"Report generated: {report_path}")


if __name__ == "__main__":
    main()
